import sys
import subprocess

try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation

def main():
    prs = Presentation(r'c:\Users\master\monsmecta-landing\public\assets\(한)리키(LiQi) 건호형 리플렛.pptx')
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                print(f'--- Slide {i+1} Table ---')
                for row in shape.table.rows:
                    print(' | '.join([cell.text.replace('\n', ' ').replace('\x0b', ' ').strip() for cell in row.cells]))
                print('---')

if __name__ == '__main__':
    main()
